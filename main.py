import os
from dotenv import load_dotenv
load_dotenv()
from fastapi import FastAPI
from fastapi.responses import FileResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
from groq import Groq

from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.enums import TA_LEFT
from reportlab.lib.units import mm

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

app = FastAPI()

app.mount("/static", StaticFiles(directory="static"), name="static")


class Pergunta(BaseModel):
    pergunta: str
    conteudo: str | None = None


client = Groq(api_key=os.getenv("GROQ_KEY"))


def gerar_plano(texto_tema: str) -> str:
    prompt = f"""
    Crie um plano de aula completo sobre: {texto_tema}

    Inclua:
    - Objetivo geral
    - Objetivos específicos
    - Conteúdo programático
    - Exemplos práticos
    - Perguntas para discussão
    - Conclusão

    Títulos devem estar em negrito como:
    **Objetivo Geral**
    **Objetivos Específicos**
    **Conteúdo Programático**
    **Exemplos Práticos**
    **Perguntas para Discussão**
    **Conclusão**

    Conteúdo em português do Brasil.
    Não escreva palavras como "Slide".
    """

    resposta = client.chat.completions.create(
    model="llama-3.1-8b-instant",
    messages=[{"role": "user", "content": prompt}]
)



    return resposta.choices[0].message.content


@app.post("/chat")
async def gerar_resposta(pergunta: Pergunta):
    try:
        texto = gerar_plano(pergunta.pergunta)
        return {"resposta": texto}
    except Exception as e:
        return JSONResponse(status_code=500, content={"erro": str(e)})


@app.post("/pdf")
async def gerar_pdf(pergunta: Pergunta):
    try:
        texto = pergunta.conteudo or gerar_plano(pergunta.pergunta)
        pdf_path = "Plano-de-Aula-UNISUL.pdf"
        margem = 18 * mm

        doc = SimpleDocTemplate(
            pdf_path,
            pagesize=A4,
            rightMargin=margem,
            leftMargin=margem,
            topMargin=margem + 30,
            bottomMargin=margem
        )

        estilo_titulo = ParagraphStyle(
            name="Titulo",
            fontSize=16,
            leading=20,
            alignment=TA_LEFT,
            spaceAfter=10,
            fontName="Helvetica-Bold"
        )

        estilo_secao = ParagraphStyle(
            name="Secao",
            fontSize=13,
            leading=18,
            alignment=TA_LEFT,
            spaceAfter=6,
            fontName="Helvetica-Bold"
        )

        estilo_conteudo = ParagraphStyle(
            name="Conteudo",
            fontSize=11,
            leading=15,
            alignment=TA_LEFT,
            spaceAfter=4
        )

        elementos = []
        elementos.append(Paragraph(f"<b>Plano de Aula — UNISUL</b>", estilo_titulo))
        elementos.append(Paragraph(f"<b>Tema:</b> {pergunta.pergunta}", estilo_conteudo))
        elementos.append(Spacer(1, 10))

        for linha in texto.split("\n"):
            linha = linha.strip()
            if not linha:
                continue

            # Títulos em markdown: **Titulo**
            if linha.startswith("**") and linha.endswith("**"):
                titulo_secao = linha.strip("* ").strip()
                elementos.append(Paragraph(f"<b>{titulo_secao}</b>", estilo_secao))
                continue

            # Itens de lista começando com "-"
            if linha.startswith("-"):
                texto_item = linha[1:].strip()
                elementos.append(Paragraph(f"- {texto_item}", estilo_conteudo))
                continue

            # Qualquer outra linha normal
            elementos.append(Paragraph(linha, estilo_conteudo))

        doc.build(elementos)
        return FileResponse(pdf_path)

    except Exception as e:
        return JSONResponse(status_code=500, content={"erro_pdf": str(e)})

@app.post("/ppt")
async def gerar_ppt(pergunta: Pergunta):
    try:
        texto = pergunta.conteudo or gerar_plano(pergunta.pergunta)

        ppt = Presentation()
        azul = RGBColor(0, 51, 153)

        temas = {
            "Objetivo Geral": "Introdução",
            "Objetivos Específicos": "Objetivos",
            "Conteúdo Programático": "Conteúdo Técnico",
            "Exemplos Práticos": "Exemplos",
            "Perguntas para Discussão": "Discussão",
            "Conclusão": "Conclusão"
        }

        secoes = {}
        chave_atual = None

        for linha in texto.split("\n"):
            linha = linha.strip().replace("*", "")

            if linha in temas:
                chave_atual = linha
                secoes[chave_atual] = []
            elif chave_atual and linha:
                secoes[chave_atual].append(linha)

        # Slide inicial
        slide = ppt.slides.add_slide(ppt.slide_layouts[6])
        titulo = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
        tf = titulo.text_frame
        tf.text = "Plano de Aula — UNISUL"
        tf.paragraphs[0].font.size = Pt(42)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = azul

        subt = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
        subt_tf = subt.text_frame
        subt_tf.text = pergunta.pergunta
        subt_tf.paragraphs[0].font.size = Pt(32)

        logo = slide.shapes.add_picture("static/unisul.png", Inches(7.5), Inches(0.3), height=Inches(1))

        for chave, conteudo in secoes.items():
            slide = ppt.slides.add_slide(ppt.slide_layouts[1])
            titulo = slide.shapes.title
            titulo.text = temas[chave]
            titulo.text_frame.paragraphs[0].font.size = Pt(40)
            titulo.text_frame.paragraphs[0].font.bold = True
            titulo.text_frame.paragraphs[0].font.color.rgb = azul

            corpo = slide.shapes.placeholders[1]
            corpo.text = ""
            for item in conteudo:
                p = corpo.text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(24)
                p.level = 0

        arquivo = "Slides-UNISUL.pptx"
        ppt.save(arquivo)
        return FileResponse(arquivo)

    except Exception as e:
        return JSONResponse(status_code=500, content={"erro_ppt": str(e)})


@app.get("/")
async def root():
    return FileResponse("templates/index.html")
